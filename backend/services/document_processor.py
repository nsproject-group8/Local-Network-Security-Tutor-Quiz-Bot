import os
from pathlib import Path
from typing import List, Dict, Any
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import markdown
from bs4 import BeautifulSoup
from loguru import logger
from config import settings

class DocumentProcessor:
    """Process various document formats for ingestion into the vector database."""
    
    def __init__(self):
        self.supported_extensions = {'.pdf', '.docx', '.pptx', '.txt', '.md'}
    
    def process_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PDF file."""
        chunks = []
        try:
            reader = PdfReader(file_path)
            filename = os.path.basename(file_path)
            
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text.strip():
                    chunks.append({
                        'text': text.strip(),
                        'metadata': {
                            'source': filename,
                            'page': page_num,
                            'type': 'pdf'
                        }
                    })
            
            logger.info(f"Extracted {len(chunks)} pages from PDF: {filename}")
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
        
        return chunks
    
    def process_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from DOCX file."""
        chunks = []
        try:
            doc = Document(file_path)
            filename = os.path.basename(file_path)
            
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            
            # Combine into larger chunks (by paragraphs)
            text = "\n\n".join(full_text)
            if text:
                chunks.append({
                    'text': text,
                    'metadata': {
                        'source': filename,
                        'type': 'docx'
                    }
                })
            
            logger.info(f"Extracted text from DOCX: {filename}")
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
        
        return chunks
    
    def process_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from PPTX file."""
        chunks = []
        try:
            prs = Presentation(file_path)
            filename = os.path.basename(file_path)
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                
                if text_parts:
                    chunks.append({
                        'text': "\n".join(text_parts),
                        'metadata': {
                            'source': filename,
                            'slide': slide_num,
                            'type': 'pptx'
                        }
                    })
            
            logger.info(f"Extracted {len(chunks)} slides from PPTX: {filename}")
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
        
        return chunks
    
    def process_text(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract text from plain text or markdown file."""
        chunks = []
        try:
            filename = os.path.basename(file_path)
            ext = os.path.splitext(file_path)[1].lower()
            
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            if text.strip():
                chunks.append({
                    'text': text.strip(),
                    'metadata': {
                        'source': filename,
                        'type': ext[1:]  # Remove the dot
                    }
                })
            
            logger.info(f"Extracted text from {ext.upper()}: {filename}")
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {e}")
        
        return chunks
    
    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        """Process a file based on its extension."""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            return self.process_pdf(file_path)
        elif ext == '.docx':
            return self.process_docx(file_path)
        elif ext == '.pptx':
            return self.process_pptx(file_path)
        elif ext in {'.txt', '.md'}:
            return self.process_text(file_path)
        else:
            logger.warning(f"Unsupported file type: {ext}")
            return []
    
    def process_directory(self, directory_path: str) -> List[Dict[str, Any]]:
        """Process all supported documents in a directory."""
        all_chunks = []
        directory = Path(directory_path)
        
        if not directory.exists():
            logger.error(f"Directory does not exist: {directory_path}")
            return all_chunks
        
        for file_path in directory.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_extensions:
                chunks = self.process_file(str(file_path))
                all_chunks.extend(chunks)
        
        logger.info(f"Processed {len(all_chunks)} total chunks from directory: {directory_path}")
        return all_chunks

# Singleton instance
document_processor = DocumentProcessor()
